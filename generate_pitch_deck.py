from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── Colour Palette ────────────────────────────────────────────────────────────
BLACK      = RGBColor(0x00, 0x00, 0x00)
NEAR_BLACK = RGBColor(0x0A, 0x0A, 0x0A)
DARK_GRAY  = RGBColor(0x18, 0x18, 0x18)
MID_GRAY   = RGBColor(0x2A, 0x2A, 0x2A)
GOLD       = RGBColor(0xF5, 0xC5, 0x18)        # main accent
GOLD_DIM   = RGBColor(0xB8, 0x8C, 0x0A)        # dimmed gold for sub-text
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xCC, 0xCC, 0xCC)
MED_GRAY   = RGBColor(0x88, 0x88, 0x88)
GREEN      = RGBColor(0x22, 0xC5, 0x5E)
RED_DIM    = RGBColor(0xEF, 0x44, 0x44)
CYAN       = RGBColor(0x06, 0xB6, 0xD4)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)


def new_prs():
    prs = Presentation()
    prs.slide_width  = SLIDE_W
    prs.slide_height = SLIDE_H
    return prs


def blank_slide(prs):
    blank_layout = prs.slide_layouts[6]   # completely blank
    return prs.slides.add_slide(blank_layout)


def fill_bg(slide, color=NEAR_BLACK):
    """Fill slide background with solid colour."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, l, t, w, h, fill_color=None, line_color=None, line_width_pt=0):
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE_TYPE.RECTANGLE
        Inches(l), Inches(t), Inches(w), Inches(h)
    )
    shape.line.fill.background() if line_color is None else None
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    else:
        shape.fill.background()
    if line_color:
        shape.line.color.rgb = line_color
        shape.line.width = Pt(line_width_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, text, l, t, w, h,
                font_size=18, bold=False, color=WHITE,
                align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = wrap
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Segoe UI"
    return txBox


def add_multiline(slide, lines, l, t, w, h,
                  font_size=14, color=LIGHT_GRAY, line_spacing_pt=None,
                  bold_first=False, first_color=None):
    """Add a textbox with multiple bullet lines."""
    txBox = slide.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    txBox.word_wrap = True
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, line in enumerate(lines):
        p = tf.add_paragraph() if i > 0 else tf.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT
        if line_spacing_pt:
            from pptx.util import Pt as _Pt
            p.space_before = _Pt(line_spacing_pt)
        run = p.add_run()
        run.text = line
        run.font.size = Pt(font_size)
        run.font.name = "Segoe UI"
        if i == 0 and bold_first:
            run.font.bold = True
            run.font.color.rgb = first_color or GOLD
        else:
            run.font.bold = False
            run.font.color.rgb = color
    return txBox


def gold_line(slide, t, l=0.5, w=12.33, height=0.03):
    """Thin gold horizontal rule."""
    add_rect(slide, l, t, w, height, fill_color=GOLD)


def slide_header(slide, number_str, title, subtitle=None):
    """Standard slide header: slide number badge + title + gold underline."""
    # Number badge
    add_rect(slide, 0.4, 0.25, 0.55, 0.5, fill_color=GOLD)
    add_textbox(slide, number_str, 0.4, 0.25, 0.55, 0.5,
                font_size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)
    # Title
    add_textbox(slide, title, 1.1, 0.22, 11.2, 0.6,
                font_size=30, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    # Subtitle
    if subtitle:
        add_textbox(slide, subtitle, 1.1, 0.82, 11.2, 0.35,
                    font_size=13, bold=False, color=MED_GRAY, align=PP_ALIGN.LEFT)
    # Gold rule
    gold_line(slide, 1.15)


# ══════════════════════════════════════════════════════════════════════════════
#  SLIDES
# ══════════════════════════════════════════════════════════════════════════════

def slide_cover(prs):
    """Cover / Title slide."""
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)

    # Gold accent bar left
    add_rect(s, 0, 0, 0.08, 7.5, fill_color=GOLD)

    # Central block
    add_textbox(s, "BlackSlon Protocol", 0.6, 1.6, 12.0, 1.0,
                font_size=52, bold=True, color=GOLD, align=PP_ALIGN.LEFT)
    add_textbox(s, "The New Dimension of European Gas & Power Wholesale Markets",
                0.6, 2.65, 10.5, 0.8,
                font_size=22, bold=False, color=WHITE, align=PP_ALIGN.LEFT)

    gold_line(s, 3.55, l=0.6, w=10.0)

    add_textbox(s, "Zero Expiry.  Zero Barriers.  100 kWh = 1 Token.",
                0.6, 3.72, 10.5, 0.55,
                font_size=17, bold=False, color=LIGHT_GRAY, align=PP_ALIGN.LEFT,
                italic=True)

    # Bottom strip
    add_rect(s, 0, 6.85, 13.33, 0.65, fill_color=MID_GRAY)
    add_textbox(s, "blackslon.org  ·  trading@blackslon.org",
                0.6, 6.88, 7.0, 0.4,
                font_size=12, color=MED_GRAY, align=PP_ALIGN.LEFT)
    add_textbox(s, "Confidential — Seed Round 2026",
                6.5, 6.88, 6.3, 0.4,
                font_size=12, color=MED_GRAY, align=PP_ALIGN.RIGHT)


def slide_problem(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "01", "The Architecture of Exclusion",
                 "PROBLEM")

    bullets = [
        ("European wholesale power & gas market: ~€500B annually, trillions in financial turnover", WHITE),
        ("Controlled by fewer than 250 institutional players", LIGHT_GRAY),
        ("Entry requires €3–5M in capital, 12+ months regulatory onboarding,\n    investment-grade credit ratings", LIGHT_GRAY),
        ("No platform exists giving anyone outside this inner circle access", LIGHT_GRAY),
    ]

    y = 1.55
    for text, col in bullets:
        # bullet dot
        add_rect(s, 0.6, y + 0.07, 0.12, 0.12, fill_color=GOLD)
        add_textbox(s, text, 0.9, y, 11.5, 0.55,
                    font_size=15, color=col)
        y += 0.72

    # Bold conclusion box
    add_rect(s, 0.5, y + 0.15, 12.3, 0.72, fill_color=MID_GRAY,
             line_color=GOLD, line_width_pt=1.5)
    add_textbox(s, "This is not a market inefficiency.  It is the architecture.",
                0.7, y + 0.22, 12.0, 0.5,
                font_size=18, bold=True, color=GOLD, align=PP_ALIGN.CENTER)

    # Right-side stat box
    add_rect(s, 9.6, 1.4, 3.2, 2.4, fill_color=DARK_GRAY,
             line_color=GOLD_DIM, line_width_pt=1)
    add_textbox(s, "< 250", 9.7, 1.5, 3.0, 0.7,
                font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
    add_textbox(s, "institutional players\ncontrol the entire\nEuropean energy market",
                9.7, 2.1, 3.0, 0.9,
                font_size=11, color=MED_GRAY, align=PP_ALIGN.CENTER)


def slide_solution(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "02", "Zero Expiry.  Zero Barriers.  100 kWh = 1 Token.",
                 "SOLUTION")

    items = [
        ("BlackSlon", "the first decentralised protocol for European wholesale gas & power markets"),
        ("Minimum entry:", "from €744,600 (smallest German yearly power contract) to a few euros"),
        ("Same market,", "accessible to any participant, anywhere, at any time"),
        ("No forced rolls,", "no expiry, no overnight costs, no clearing banks"),
    ]

    y = 1.55
    for label, desc in items:
        add_textbox(s, label, 0.6, y, 3.2, 0.45,
                    font_size=16, bold=True, color=GOLD)
        add_textbox(s, desc, 3.7, y, 8.8, 0.5,
                    font_size=15, color=LIGHT_GRAY)
        gold_line(s, y + 0.55, l=0.6, w=12.2, height=0.012)
        y += 0.78

    # Big tagline bottom
    add_textbox(s, "The same market.  For everyone.",
                0.6, 6.45, 12.5, 0.65,
                font_size=22, bold=True, color=GOLD, align=PP_ALIGN.CENTER,
                italic=True)


def slide_product(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "03", "The Protocol", "PRODUCT")

    components = [
        ("BSSZ", "Physically-anchored price corridor filtering chaos, preserving trend"),
        ("BSEI", "Manipulation-resistant 72h rolling benchmark index"),
        ("Perpetual Instruments", "No forced rolls, no expiry, no overnight costs"),
        ("Smart Liquidation", "10% increments — never a full forced close"),
        ("eEURO + €BSR", "MiCA-compliant settlement & collateral layer"),
    ]

    for i, (label, desc) in enumerate(components):
        col = 0 if i < 3 else 1
        row = i if i < 3 else i - 3
        lx = 0.5 + col * 6.5
        ty = 1.55 + row * 1.35

        add_rect(s, lx, ty, 6.1, 1.15, fill_color=DARK_GRAY,
                 line_color=GOLD_DIM, line_width_pt=1)
        add_textbox(s, label, lx + 0.18, ty + 0.1, 5.7, 0.38,
                    font_size=15, bold=True, color=GOLD)
        add_textbox(s, desc, lx + 0.18, ty + 0.48, 5.7, 0.55,
                    font_size=12, color=LIGHT_GRAY)

    # Demo link
    add_rect(s, 0.5, 6.5, 12.3, 0.65, fill_color=MID_GRAY)
    add_textbox(s, "Demo live:  blackslon.org/markets",
                0.7, 6.55, 12.0, 0.45,
                font_size=15, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_market(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "04", "The Opportunity", "MARKET")

    # Big number boxes
    boxes = [
        ("~€500B", "physical market\nannually"),
        ("Trillions", "financial turnover\nannually"),
        ("0", "direct decentralised\ncompetitors"),
    ]
    bx = 0.5
    for val, label in boxes:
        add_rect(s, bx, 1.5, 3.9, 1.8, fill_color=DARK_GRAY,
                 line_color=GOLD_DIM, line_width_pt=1)
        add_textbox(s, val, bx + 0.1, 1.6, 3.7, 0.75,
                    font_size=36, bold=True, color=GOLD, align=PP_ALIGN.CENTER)
        add_textbox(s, label, bx + 0.1, 2.3, 3.7, 0.75,
                    font_size=12, color=MED_GRAY, align=PP_ALIGN.CENTER)
        bx += 4.2

    # Currently excluded
    add_textbox(s, "Currently excluded from these markets:", 0.5, 3.55, 12.0, 0.4,
                font_size=14, bold=True, color=GOLD)
    excl = [
        "Independent traders & proprietary shops",
        "Industrial consumers seeking direct hedges",
        "Renewable energy producers",
        "Retail & institutional investors outside the club",
    ]
    y = 4.05
    for item in excl:
        add_rect(s, 0.55, y + 0.05, 0.1, 0.1, fill_color=GOLD)
        add_textbox(s, item, 0.8, y, 11.5, 0.42,
                    font_size=13, color=LIGHT_GRAY)
        y += 0.5

    # Adjacent note
    add_textbox(s,
                "Adjacent: crypto derivatives $50B+ daily volume  ·  commodity tokenisation (emerging)",
                0.5, 6.55, 12.5, 0.55,
                font_size=11, color=MED_GRAY, align=PP_ALIGN.LEFT, italic=True)


def slide_bizmodel(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "05", "How BlackSlon Makes Money", "BUSINESS MODEL")

    streams = [
        ("Trading Fees", "0.20% – 1.00% per transaction\n(tiered by €BSR collateral level)"),
        ("Maintenance Fee", "0.1% / month on total vault value"),
        ("PLP Fees", "Proportional share from liquidity\nprovision pool"),
        ("€BSR Appreciation", "Deflationary token — burns on\nliquidation events"),
    ]

    sx, sy = 0.5, 1.55
    for i, (stream, detail) in enumerate(streams):
        col = i % 2
        row = i // 2
        lx = sx + col * 6.4
        ty = sy + row * 2.1
        add_rect(s, lx, ty, 6.0, 1.85, fill_color=DARK_GRAY,
                 line_color=GOLD, line_width_pt=1.5)
        add_textbox(s, stream, lx + 0.2, ty + 0.12, 5.6, 0.45,
                    font_size=16, bold=True, color=GOLD)
        add_textbox(s, detail, lx + 0.2, ty + 0.58, 5.6, 0.95,
                    font_size=13, color=LIGHT_GRAY)

    # Asset-light tag
    add_textbox(s,
                "Asset-light model — protocol runs algorithmically, 24 / 7 / 365",
                0.5, 6.48, 12.5, 0.55,
                font_size=14, bold=True, color=MED_GRAY, align=PP_ALIGN.CENTER,
                italic=True)


def slide_traction(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "06", "Where We Are", "TRACTION")

    milestones = [
        (True,  "White Paper v3.0",
                "Complete protocol specification — BSSZ, BSEI, ADR, Health Factor, risk management"),
        (True,  "Working Demo",
                "blackslon.org/markets — live order book, matching engine, risk UI"),
        (True,  "Email Wallet Onboarding",
                "OTP-verified wallet connection — live"),
        (True,  "Alliance DAO Application",
                "Accelerator application submitted"),
        (False, "CASP Licensing",
                "Lithuania / Bulgaria — Q3 2026"),
        (False, "Smart Contract Audit",
                "EVM deployment — H2 2026"),
    ]

    y = 1.5
    for done, title, desc in milestones:
        # Status dot
        dot_color = GREEN if done else MED_GRAY
        add_rect(s, 0.5, y + 0.12, 0.18, 0.18, fill_color=dot_color)
        add_textbox(s, "✓" if done else "○", 0.5, y + 0.08, 0.22, 0.28,
                    font_size=10, bold=True, color=BLACK if done else MED_GRAY,
                    align=PP_ALIGN.CENTER)
        add_textbox(s, title, 0.85, y, 4.2, 0.38,
                    font_size=14, bold=True, color=GOLD if done else MED_GRAY)
        add_textbox(s, desc, 5.2, y, 7.8, 0.38,
                    font_size=12, color=LIGHT_GRAY if done else MED_GRAY)
        gold_line(s, y + 0.48, l=0.5, w=12.3, height=0.008)
        y += 0.72

    add_textbox(s, "Founder network: 20+ years of relationships across ICE, EEX, TGE, EPEX",
                0.5, 6.55, 12.5, 0.55,
                font_size=12, color=GOLD_DIM, italic=True)


def slide_competition(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "07", "Why No One Has Done This", "COMPETITION")

    headers = ["", "BlackSlon", "Traditional\nExchanges", "DeFi\nProtocols"]
    rows = [
        ("European energy",       True,  True,  False),
        ("Open access",           True,  False, True),
        ("Physical anchor",       True,  True,  False),
        ("Perpetual instruments", True,  False, None),
        ("No expiry / rollover",  True,  False, False),
    ]

    col_widths = [3.6, 2.4, 2.4, 2.4]
    col_starts = [0.5, 4.2, 6.7, 9.2]
    row_h = 0.72
    table_top = 1.45

    # Header row
    for ci, (hdr, cx) in enumerate(zip(headers, col_starts)):
        bg = GOLD if ci == 1 else MID_GRAY
        fg = BLACK if ci == 1 else MED_GRAY
        add_rect(s, cx, table_top, col_widths[ci], 0.58, fill_color=bg)
        add_textbox(s, hdr, cx + 0.05, table_top + 0.05,
                    col_widths[ci] - 0.1, 0.5,
                    font_size=12, bold=True, color=fg, align=PP_ALIGN.CENTER)

    # Data rows
    for ri, (label, *vals) in enumerate(rows):
        ty = table_top + 0.58 + ri * row_h
        row_bg = DARK_GRAY if ri % 2 == 0 else MID_GRAY

        for ci, (cx, val) in enumerate(zip(col_starts, [label] + list(vals))):
            add_rect(s, cx, ty, col_widths[ci], row_h - 0.04,
                     fill_color=GOLD if ci == 1 else row_bg)
            if ci == 0:
                add_textbox(s, val, cx + 0.15, ty + 0.18,
                            col_widths[ci] - 0.2, 0.38,
                            font_size=12, color=LIGHT_GRAY)
            else:
                if val is True:
                    sym, col = "YES", GREEN if ci != 1 else BLACK
                elif val is False:
                    sym, col = "NO", RED_DIM if ci != 1 else BLACK
                else:
                    sym, col = "Partial", GOLD_DIM if ci != 1 else BLACK
                add_textbox(s, sym, cx, ty + 0.18,
                            col_widths[ci], 0.38,
                            font_size=13, bold=(ci == 1), color=col,
                            align=PP_ALIGN.CENTER)

    add_textbox(s,
                "No direct competitor.  Closest: energy derivatives on CME/ICE (institutional only) "
                "and synthetic commodity protocols (no physical anchor).",
                0.5, 6.42, 12.5, 0.7,
                font_size=11, color=MED_GRAY, italic=True)


def slide_team(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "08", "The Founder", "TEAM")

    # Name block
    add_rect(s, 0.5, 1.45, 12.3, 0.65, fill_color=MID_GRAY)
    add_textbox(s, "Konrad Dynkiewicz  —  Founder",
                0.7, 1.52, 12.0, 0.48,
                font_size=22, bold=True, color=GOLD)

    achievements = [
        "20+ years trading Natural Gas, Power, Oil Products, Carbon Emissions across European markets",
        "Trader → Originator → Director → Partner in state-owned, private and global trading houses",
        "One of the largest natural gas suppliers from the European Union to Ukraine",
        "Operated in CEE frontier markets before digital benchmarks or centralised exchanges existed",
        "University of Warsaw (UW)  ·  St. Petersburg State University  ·  MGIMO Moscow",
    ]

    y = 2.32
    for item in achievements:
        add_rect(s, 0.6, y + 0.09, 0.1, 0.1, fill_color=GOLD)
        add_textbox(s, item, 0.88, y, 11.6, 0.46, font_size=13, color=LIGHT_GRAY)
        y += 0.62

    # Seeking box
    add_rect(s, 0.5, 6.3, 12.3, 0.82, fill_color=DARK_GRAY,
             line_color=GOLD, line_width_pt=1.5)
    add_textbox(s, "Seeking:  Technical co-founder (blockchain / smart contracts)  +  Regulatory counsel",
                0.7, 6.42, 12.0, 0.55,
                font_size=13, bold=True, color=GOLD, align=PP_ALIGN.CENTER)


def slide_roadmap(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "09", "The Path", "ROADMAP")

    phases = [
        ("Phase 1", "2026", "Virtual Protocol",
         "CASP licensing · Smart contract audit\nOracle infrastructure · Live trading · PLP onboarding",
         GOLD),
        ("Phase 2", "2027–2028", "Physical Integration",
         "BS-P/G token physical redemption · Industrial consumer hedging\nCross-market swaps · RWA bridge",
         CYAN),
        ("Phase 3", "2028+", "Benchmark",
         "BSEI becomes European energy reference\nExpansion to 15+ markets",
         GREEN),
    ]

    px = 0.5
    for phase, year, title, detail, accent in phases:
        add_rect(s, px, 1.45, 4.1, 4.5, fill_color=DARK_GRAY,
                 line_color=accent, line_width_pt=2)
        # Phase label
        add_rect(s, px, 1.45, 4.1, 0.55, fill_color=accent)
        add_textbox(s, phase, px + 0.1, 1.48, 2.0, 0.45,
                    font_size=14, bold=True, color=BLACK)
        add_textbox(s, year, px + 2.1, 1.48, 1.9, 0.45,
                    font_size=14, bold=True, color=BLACK, align=PP_ALIGN.RIGHT)
        # Title
        add_textbox(s, title, px + 0.15, 2.12, 3.8, 0.55,
                    font_size=17, bold=True, color=accent)
        # Detail
        add_textbox(s, detail, px + 0.15, 2.75, 3.8, 2.8,
                    font_size=12, color=LIGHT_GRAY)
        px += 4.3

    # Arrow connectors (simple triangles / text)
    add_textbox(s, "▶", 4.35, 3.55, 0.45, 0.5,
                font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)
    add_textbox(s, "▶", 8.65, 3.55, 0.45, 0.5,
                font_size=20, color=MID_GRAY, align=PP_ALIGN.CENTER)

    add_textbox(s,
                "All timelines subject to regulatory approvals and market conditions.",
                0.5, 6.5, 12.5, 0.55,
                font_size=10, color=MED_GRAY, italic=True)


def slide_ask(prs):
    s = blank_slide(prs)
    fill_bg(s, NEAR_BLACK)
    slide_header(s, "10", "Join the Build", "INVESTMENT ASK")

    # Raise box
    add_rect(s, 0.5, 1.42, 5.8, 1.1, fill_color=GOLD)
    add_textbox(s, "Raising  €2–5M  Seed", 0.65, 1.52, 5.5, 0.75,
                font_size=26, bold=True, color=BLACK, align=PP_ALIGN.CENTER)

    # Allocation
    add_textbox(s, "Allocation", 0.5, 2.72, 5.8, 0.4,
                font_size=14, bold=True, color=GOLD)
    alloc = [
        ("CASP licensing (Lithuania / Bulgaria)", "25%"),
        ("Smart contract development & audit", "30%"),
        ("Oracle infrastructure", "15%"),
        ("Regulatory counsel", "15%"),
        ("Operations & runway (12 months)", "15%"),
    ]
    y = 3.18
    for item, pct in alloc:
        add_textbox(s, item, 0.6, y, 4.3, 0.38, font_size=12, color=LIGHT_GRAY)
        add_textbox(s, pct, 5.0, y, 1.1, 0.38,
                    font_size=13, bold=True, color=GOLD, align=PP_ALIGN.RIGHT)
        y += 0.46

    # Looking for (right column)
    add_rect(s, 6.8, 1.42, 6.0, 4.75, fill_color=DARK_GRAY,
             line_color=GOLD_DIM, line_width_pt=1)
    add_textbox(s, "Looking for:", 7.0, 1.55, 5.6, 0.42,
                font_size=14, bold=True, color=GOLD)
    looking = [
        ("Seed capital", "VC / angel"),
        ("Technical co-founder", "blockchain / smart contracts partner"),
        ("First PLP", "€500K minimum · Genesis €BSR allocation"),
    ]
    ly = 2.08
    for title, sub in looking:
        add_rect(s, 6.95, ly, 5.65, 1.0, fill_color=MID_GRAY)
        add_textbox(s, title, 7.1, ly + 0.08, 5.4, 0.38,
                    font_size=14, bold=True, color=WHITE)
        add_textbox(s, sub, 7.1, ly + 0.5, 5.4, 0.35,
                    font_size=11, color=MED_GRAY)
        ly += 1.18

    # Contact footer
    add_rect(s, 0, 6.72, 13.33, 0.78, fill_color=GOLD)
    add_textbox(s, "trading@blackslon.org    ·    blackslon.org    ·    blackslon.org/markets",
                0.5, 6.82, 12.5, 0.5,
                font_size=16, bold=True, color=BLACK, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
#  BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build():
    prs = new_prs()
    slide_cover(prs)
    slide_problem(prs)
    slide_solution(prs)
    slide_product(prs)
    slide_market(prs)
    slide_bizmodel(prs)
    slide_traction(prs)
    slide_competition(prs)
    slide_team(prs)
    slide_roadmap(prs)
    slide_ask(prs)

    out = "BlackSlon_PitchDeck_2026.pptx"
    prs.save(out)
    print(f"Saved: {out}")


if __name__ == "__main__":
    build()
